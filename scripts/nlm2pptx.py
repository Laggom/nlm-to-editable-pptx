"""
nlm2pptx — NotebookLM(이미지 기반) PPTX/PDF를 '편집 가능한' PPTX로 변환하는 핵심 로직.

설계 원칙 (이식성):
- 표준 라이브러리 + python-pptx + pillow 만 사용. OpenAI 호출은 urllib(표준)로 직접.
- 경로/환경 의존 없음: 모든 입출력은 함수 인자. OS 무관(Windows/Linux/macOS/Databricks).
- API 키는 환경변수 OPENAI_API_KEY 에서 읽음 (인자로 직접 주입도 가능).
- macOS 전용 도구(editppt, open, 시스템 폰트 경로) 일절 사용 안 함.

파이프라인:
  1) extract_slides(): pptx/pdf → 페이지별 PNG
  2) erase_text():     각 PNG → 글자 지운 배경 PNG  (gpt-image-2, images/edits)
  3) ocr_slide():      각 PNG → 텍스트 블록 JSON     (gpt-5.5, vision + JSON)
  4) build_pptx():     배경 + OCR블록 → 편집 가능 PPTX (python-pptx)

CLI:
  python nlm2pptx.py <input.pptx|input.pdf> <output.pptx> [--workdir DIR] [--no-erase] [--font "맑은 고딕"]
"""
from __future__ import annotations
import base64, json, os, re, sys, time, uuid, zipfile, io, glob, argparse, logging, threading
import urllib.request, urllib.error
from concurrent.futures import ThreadPoolExecutor, as_completed

__version__ = "1.1.0"

# ──────────────────────────────────────────────────────────────────────────
# .env 로더 (표준 라이브러리만 사용). CWD / 스크립트 디렉토리 / 그 상위에서
# .env 를 찾아 KEY=VALUE 를 os.environ 에 채운다(기존 환경변수는 덮어쓰지 않음).
# ──────────────────────────────────────────────────────────────────────────
def _load_dotenv() -> None:
    here = os.path.dirname(os.path.abspath(__file__))
    candidates = [
        os.path.join(os.getcwd(), ".env"),
        os.path.join(here, ".env"),
        os.path.join(os.path.dirname(here), ".env"),
    ]
    seen = set()
    for path in candidates:
        if path in seen or not os.path.isfile(path):
            continue
        seen.add(path)
        try:
            with open(path, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line or line.startswith("#") or "=" not in line:
                        continue
                    if line.startswith("export "):
                        line = line[len("export "):]
                    name, _, value = line.partition("=")
                    name = name.strip()
                    value = value.strip().strip('"').strip("'")
                    if name and name not in os.environ:
                        os.environ[name] = value
        except OSError:
            pass

_load_dotenv()

# ──────────────────────────────────────────────────────────────────────────
# 로깅: 타임스탬프 + 단계/슬라이드/소요시간. 콘솔 + (workdir 지정 시) 파일.
# 환경변수 NLM2PPTX_LOG_LEVEL 로 조절(기본 INFO, 디버그는 DEBUG).
# ──────────────────────────────────────────────────────────────────────────
log = logging.getLogger("nlm2pptx")
if not log.handlers:
    _h = logging.StreamHandler(sys.stdout)
    _h.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(message)s", "%H:%M:%S"))
    log.addHandler(_h)
log.setLevel(getattr(logging, os.environ.get("NLM2PPTX_LOG_LEVEL", "INFO").upper(), logging.INFO))

def _add_file_log(workdir: str):
    """workdir/convert.log 에 전체 로그를 함께 기록."""
    fh = logging.FileHandler(os.path.join(workdir, "convert.log"), encoding="utf-8")
    fh.setFormatter(logging.Formatter("%(asctime)s [%(levelname)s] %(message)s"))
    fh.setLevel(logging.DEBUG)
    log.addHandler(fh)
    return fh

# ──────────────────────────────────────────────────────────────────────────
# 설정 (기본 모델)
# ──────────────────────────────────────────────────────────────────────────
IMAGE_MODEL = os.environ.get("NLM2PPTX_IMAGE_MODEL", "gpt-image-2")   # 글자 제거
OCR_MODEL   = os.environ.get("NLM2PPTX_OCR_MODEL", "gpt-5.5")          # OCR
DEFAULT_FONT = os.environ.get("NLM2PPTX_FONT", "맑은 고딕")            # 한글 폰트(이름만 기입)
OPENAI_BASE = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1")
IMAGE_SIZE  = "1536x1024"   # gpt-image 허용 가로 사이즈(3:2). 원본 비율은 레터박스로 보존(erase_text 참고).
HTTP_TIMEOUT = int(os.environ.get("NLM2PPTX_HTTP_TIMEOUT", "400"))   # API 응답 대기(초). OCR 지연 대응


def _api_key(explicit: str | None = None) -> str:
    key = explicit or os.environ.get("OPENAI_API_KEY")
    if not key:
        raise RuntimeError(
            "OpenAI API 키가 없습니다. 환경변수 OPENAI_API_KEY 를 설정하거나 api_key 인자로 전달하세요."
        )
    return key


# ──────────────────────────────────────────────────────────────────────────
# 1) 슬라이드 이미지 추출
# ──────────────────────────────────────────────────────────────────────────
def extract_slides(input_path: str, out_dir: str) -> list[str]:
    """pptx 또는 pdf에서 페이지별 PNG를 out_dir에 추출. 정렬된 경로 리스트 반환."""
    os.makedirs(out_dir, exist_ok=True)
    ext = os.path.splitext(input_path)[1].lower()
    paths: list[str] = []

    if ext in (".pptx", ".ppt"):
        # pptx는 zip. ppt/media/imageN.* 가 슬라이드 풀페이지 이미지.
        with zipfile.ZipFile(input_path) as z:
            media = sorted(
                [n for n in z.namelist() if n.startswith("ppt/media/image")],
                key=lambda n: int(re.search(r"image(\d+)", n).group(1)),
            )
            for i, name in enumerate(media, 1):
                data = z.read(name)
                # PNG로 정규화 (pillow)
                from PIL import Image
                im = Image.open(io.BytesIO(data)).convert("RGB")
                p = os.path.join(out_dir, f"slide{i}.png")
                im.save(p)
                paths.append(p)

    elif ext == ".pdf":
        # PDF는 PyMuPDF(fitz)로 페이지 렌더. (없으면 안내 후 예외)
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise RuntimeError("PDF 입력은 PyMuPDF가 필요합니다: pip install pymupdf")
        doc = fitz.open(input_path)
        for i, page in enumerate(doc, 1):
            pix = page.get_pixmap(dpi=200)
            p = os.path.join(out_dir, f"slide{i}.png")
            pix.save(p)
            paths.append(p)
    else:
        raise ValueError(f"지원하지 않는 입력 형식: {ext} (pptx 또는 pdf)")

    if not paths:
        raise RuntimeError(f"{input_path} 에서 슬라이드 이미지를 찾지 못했습니다.")
    return paths


# ──────────────────────────────────────────────────────────────────────────
# 2) 글자 제거 (gpt-image-2, images/edits 멀티파트)
# ──────────────────────────────────────────────────────────────────────────
# composition-lock 프롬프트: "글자만 지운다"가 아니라 "비텍스트 픽셀은 손대지 마라"를 강하게 지시.
# 단순 "remove text"보다 글자 잔존/엉뚱한 내용 생성(hallucinate)이 줄고, frame 내 재구성을 억제한다.
ERASE_PROMPT = (
    "Remove ONLY the text and letters from this slide image. This is a PIXEL-PRESERVATION task, "
    "NOT a creative regeneration. Keep the EXACT same composition: do NOT move, resize, shift, "
    "crop, recolor, or redraw ANY non-text element. Every graph, axis, curve, diagram, "
    "illustration, box, line, gridline, and background region MUST stay at the EXACT same "
    "position and scale as the input — only the glyphs disappear, replaced by the background "
    "that was behind them. Do NOT add any new text, marks, or labels, and do NOT invent content "
    "in the cleared spots; fill them with the surrounding background only."
)

def _letterbox(im, tw: int, th: int):
    """im을 tw×th 캔버스에 비율 유지(레터박스)로 배치. 반환 (padded, (left,top,w,h))."""
    W, H = im.size
    if W / H > tw / th:                 # 원본이 더 넓음 → 폭 맞추고 상하 여백
        sw, sh = tw, max(1, round(tw * H / W)); left, top = 0, (th - sh) // 2
    else:                               # 더 좁음 → 높이 맞추고 좌우 여백
        sh, sw = th, max(1, round(th * W / H)); top, left = 0, (tw - sw) // 2
    from PIL import Image
    base = Image.new("RGB", (tw, th), (255, 255, 255))
    base.paste(im.resize((sw, sh)), (left, top))
    return base, (left, top, sw, sh)

def erase_text(png_path: str, out_path: str, api_key: str | None = None,
               model: str | None = None, retries: int = 3) -> str:
    """슬라이드 PNG에서 글자만 제거한 배경 이미지를 생성해 out_path에 저장.

    종횡비 보존: gpt-image는 출력 사이즈가 고정(IMAGE_SIZE)이라, 원본을 그대로 보내면
    비율이 다를 때 모델이 내용을 재구성(reflow)해 레이아웃이 어긋난다. 그래서
    (1) 원본을 IMAGE_SIZE 캔버스에 레터박스로 패딩해 전송 → 모델은 비율을 바꿀 필요가 없어
    재구성하지 않고 글자만 지움 → (2) 받은 결과에서 패딩 영역을 크롭해 원본 비율을 복원한다.
    """
    from PIL import Image
    key = _api_key(api_key)
    model = model or IMAGE_MODEL
    tw, th = (int(v) for v in IMAGE_SIZE.lower().split("x"))
    orig = Image.open(png_path).convert("RGB")
    padded, (left, top, sw, sh) = _letterbox(orig, tw, th)
    buf = io.BytesIO(); padded.save(buf, "PNG"); img = buf.getvalue()
    boundary = "----" + uuid.uuid4().hex
    def part_field(name, val):
        return (f'--{boundary}\r\nContent-Disposition: form-data; name="{name}"\r\n\r\n{val}\r\n').encode()
    def part_file(name, fn, data):
        return (f'--{boundary}\r\nContent-Disposition: form-data; name="{name}"; '
                f'filename="{fn}"\r\nContent-Type: image/png\r\n\r\n').encode() + data + b"\r\n"
    body = (part_field("model", model) + part_field("prompt", ERASE_PROMPT)
            + part_field("size", IMAGE_SIZE)
            + part_file("image", os.path.basename(png_path), img)
            + (f"--{boundary}--\r\n").encode())
    name = os.path.basename(png_path)
    last = None
    for a in range(retries):
        t0 = time.time()
        try:
            log.debug(f"erase {name}: 시도 {a+1}/{retries} ({model})")
            req = urllib.request.Request(
                f"{OPENAI_BASE}/images/edits", data=body,
                headers={"Authorization": f"Bearer {key}",
                         "Content-Type": f"multipart/form-data; boundary={boundary}"})
            d = json.loads(urllib.request.urlopen(req, timeout=HTTP_TIMEOUT).read())
            b64 = d["data"][0]["b64_json"]
            res = Image.open(io.BytesIO(base64.b64decode(b64))).convert("RGB")
            if res.size != (tw, th):                 # 안전망: 결과가 요청 사이즈와 다르면 맞춤
                res = res.resize((tw, th))
            res.crop((left, top, left + sw, top + sh)).save(out_path)   # 패딩 크롭 → 원본 비율 복원
            log.info(f"erase {name}: OK ({time.time()-t0:.1f}s)")
            return out_path
        except urllib.error.HTTPError as e:
            last = f"HTTP {e.code}: {e.read()[:200].decode('utf-8','ignore')}"
            log.warning(f"erase {name}: 시도 {a+1} 실패 ({time.time()-t0:.1f}s) — {last}")
            time.sleep(5)
        except Exception as e:
            last = f"{type(e).__name__}: {e}"
            log.warning(f"erase {name}: 시도 {a+1} 실패 ({time.time()-t0:.1f}s) — {last}")
            time.sleep(5)
    raise RuntimeError(f"글자 제거 실패({name}) {retries}회: {last}")


# ──────────────────────────────────────────────────────────────────────────
# 3) OCR (gpt-5.5, vision + 구조화 JSON)
# ──────────────────────────────────────────────────────────────────────────
OCR_PROMPT = """You are an expert slide OCR system. Analyze this 16:9 slide and return ONLY JSON: {"blocks":[...]}.
Each block = one visually-coherent text line or label. Fields:
- text: exact text. For math, convert to readable plain unicode (use ², ₂, √, Σ, ±, ÷, →, Greek letters; NOT LaTeX backslashes).
- box_2d: [ymin,xmin,ymax,xmax] in 0-1000 coords, TIGHT around the glyphs.
- font_size_pt: REAL point size for a 7.5-inch-tall slide (measure glyph height precisely; titles ~36-44, body ~14-20, captions ~10-13).
- bold: true/false
- align: left|center|right
- color: hex like FFFFFF
Group multi-line paragraphs into one block with \\n. Be precise about font_size_pt — it must match visual glyph height."""

def ocr_slide(png_path: str, api_key: str | None = None,
              model: str | None = None, retries: int = 3) -> list[dict]:
    """슬라이드 PNG에서 텍스트 블록 리스트(JSON) 추출."""
    key = _api_key(api_key)
    model = model or OCR_MODEL
    img_b64 = base64.b64encode(open(png_path, "rb").read()).decode()
    body = {
        "model": model,
        "messages": [{"role": "user", "content": [
            {"type": "text", "text": OCR_PROMPT},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
        ]}],
        "response_format": {"type": "json_object"},
    }
    name = os.path.basename(png_path)
    last = None
    for a in range(retries):
        t0 = time.time()
        try:
            log.debug(f"ocr {name}: 시도 {a+1}/{retries} ({model})")
            req = urllib.request.Request(
                f"{OPENAI_BASE}/chat/completions",
                data=json.dumps(body).encode(),
                headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"})
            d = json.loads(urllib.request.urlopen(req, timeout=HTTP_TIMEOUT).read())
            content = d["choices"][0]["message"]["content"]
            obj = json.loads(content)
            blocks = obj.get("blocks", obj if isinstance(obj, list) else [])
            log.info(f"ocr {name}: OK {len(blocks)}블록 ({time.time()-t0:.1f}s)")
            return blocks
        except urllib.error.HTTPError as e:
            last = f"HTTP {e.code}: {e.read()[:200].decode('utf-8','ignore')}"
            log.warning(f"ocr {name}: 시도 {a+1} 실패 ({time.time()-t0:.1f}s) — {last}")
            time.sleep(5)
        except Exception as e:
            last = f"{type(e).__name__}: {e}"
            log.warning(f"ocr {name}: 시도 {a+1} 실패 ({time.time()-t0:.1f}s) — {last}")
            time.sleep(5)
    raise RuntimeError(f"OCR 실패({name}) {retries}회: {last}")


# ──────────────────────────────────────────────────────────────────────────
# LaTeX → 유니코드 평문 (gpt-5.5가 대부분 평문으로 주지만, 혹시 모를 폴백)
# ──────────────────────────────────────────────────────────────────────────
_GREEK = {r"\\Psi":"Ψ",r"\\psi":"ψ",r"\\alpha":"α",r"\\beta":"β",r"\\gamma":"γ",r"\\delta":"δ",
          r"\\Delta":"Δ",r"\\theta":"θ",r"\\omega":"ω",r"\\Omega":"Ω",r"\\lambda":"λ",r"\\mu":"μ",
          r"\\nu":"ν",r"\\pi":"π",r"\\sigma":"σ",r"\\phi":"φ",r"\\rho":"ρ",r"\\tau":"τ",r"\\epsilon":"ε"}
_SYM = {r"\\sum":"Σ",r"\\int":"∫",r"\\partial":"∂",r"\\cdot":"·",r"\\times":"×",r"\\pm":"±",
        r"\\mp":"∓",r"\\infty":"∞",r"\\approx":"≈",r"\\neq":"≠",r"\\leq":"≤",r"\\geq":"≥",
        r"\\rightarrow":"→",r"\\to":"→",r"\\hbar":"ℏ",r"\\,":" ",r"\\!":"",r"\\quad":"  ",r"\\ ":" "}
_SUP = str.maketrans("0123456789+-=()n", "⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁼⁽⁾ⁿ")
_SUB = str.maketrans("0123456789+-=()", "₀₁₂₃₄₅₆₇₈₉₊₋₌₍₎")

def latex2plain(t: str) -> str:
    if not t:
        return t
    s = t
    for _ in range(8):
        before = s
        s = re.sub(r"\\sqrt\s*\{([^{}]*)\}", r"√(\1)", s)
        s = re.sub(r"\\frac\s*\{([^{}]*)\}\s*\{([^{}]*)\}", r"(\1)/(\2)", s)
        if s == before:
            break
    s = re.sub(r"\\sqrt\s*\{([^{}]*)\}", r"√(\1)", s)
    s = re.sub(r"\\vec\s*\{([^{}]*)\}", r"\1⃗", s)
    for k, v in {**_GREEK, **_SYM}.items():
        s = re.sub(k + r"(?![A-Za-z])", v, s)
    sup = lambda m: m.group(1).translate(_SUP) if all(c in "0123456789+-=()n" for c in m.group(1)) else "^"+m.group(1)
    sub = lambda m: m.group(1).translate(_SUB) if all(c in "0123456789+-=()" for c in m.group(1)) else "_"+m.group(1)
    s = re.sub(r"\^\{([^{}]*)\}", sup, s)
    s = re.sub(r"\^(-?[0-9]+)", sup, s)
    s = re.sub(r"\^([A-Za-z0-9])", sup, s)
    s = re.sub(r"_\{([^{}]*)\}", sub, s)
    s = s.replace("{", "").replace("}", "")
    s = re.sub(r"\\([A-Za-z]+)", r"\1", s)
    return re.sub(r"\s+", " ", s).strip()


# ──────────────────────────────────────────────────────────────────────────
# 4) PPTX 조립
# ──────────────────────────────────────────────────────────────────────────
_GLYPH_TO_FONT = 1.30     # 타이트 글리프박스 높이 → 폰트 pt 환산 계수(글자가 행 높이를 채우도록)
_FONT_SAFETY = 0.98       # 글꼴 메트릭 오차 흡수(타이트 박스에서 강제 줄바꿈 방지)

def _is_numeric_cell(t: str) -> bool:
    """숫자/기호 위주의 표 셀인지 판정(우측정렬 대상).
    OCR이 표 숫자의 정렬을 center/right로 들쭉날쭉 주므로 우측정렬로 통일하기 위함.
    예: 19,671.3 / +4,540.8 / 190% / -1%p / 110.1"""
    t = t.strip()
    if not any(c.isdigit() for c in t):
        return False
    return len(t) <= 16 and all(c.isdigit() or c in " ,.%+-()원pP△▲▽▼" for c in t)

def _em_width(s: str) -> float:
    """문자열의 em 폭 합. CJK=1.0, 공백=0.4, 그 외=0.5."""
    w = 0.0
    for ch in s:
        o = ord(ch)
        if (0x1100 <= o <= 0x11FF or 0x2E80 <= o <= 0x9FFF or 0xAC00 <= o <= 0xD7A3 or
                0xF900 <= o <= 0xFAFF or 0xFF00 <= o <= 0xFFEF or 0x3000 <= o <= 0x30FF):
            w += 1.0
        elif ch == " ":
            w += 0.33
        else:
            w += 0.55         # 라틴/숫자(맑은 고딕 실측 근사)
    return w

def _wrap_intent(text: str) -> bool:
    """줄바꿈을 허용할 '긴 문단'인지. 짧은 한 줄(라벨/숫자/제목)은 줄바꿈하지 않는다."""
    return ("\n" in text) or (sum(_em_width(ln) for ln in text.split("\n")) > 16)

def _fit_font(text: str, box_w_pt: float, box_h_pt: float) -> float:
    """OCR 추정 대신 박스 기하로 폰트 pt 산정.
    한 줄에 맞으면 박스 높이 기준, 넘치면 면적(폭×높이) 기준 → OCR 폰트 오차 자동 보정."""
    lines = text.split("\n")
    nlines = len(lines)
    longest_em = max((_em_width(ln) for ln in lines), default=1.0) or 1.0
    total_em = sum(_em_width(ln) for ln in lines) or 1.0
    f_height = (box_h_pt / nlines) * _GLYPH_TO_FONT
    if not _wrap_intent(text):
        # 짧은 한 줄(라벨/숫자/제목): 박스 높이 기준 → 행 높이 같으면 크기도 균일.
        # 폭의 1.1배까지만 허용(과도한 가로 오버플로 방지). 줄바꿈은 호출부에서 끈다.
        f = min(f_height, box_w_pt / longest_em * 1.10)
    elif longest_em * f_height <= box_w_pt:        # 여러 줄이지만 자연 크기로 들어감
        f = f_height
    else:                                          # 긴 문단 → 줄바꿈 후 높이 ≈ 박스 높이
        f = (box_w_pt * box_h_pt / (1.2 * total_em)) ** 0.5
    return max(8.0, min(f, 46.0))

def build_pptx(bg_paths: list[str], ocr_blocks_per_slide: list[list[dict]],
               out_path: str, font: str | None = None) -> str:
    """배경 이미지 + 슬라이드별 OCR 블록 → 편집 가능 PPTX 저장."""
    from pptx import Presentation
    from pptx.util import Emu, Pt, Inches
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    font = font or DEFAULT_FONT
    EMU_W, EMU_H = Inches(13.333), Inches(7.5)
    ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H
    blank = prs.slide_layouts[6]

    for bg, blocks in zip(bg_paths, ocr_blocks_per_slide):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(bg, 0, 0, width=EMU_W, height=EMU_H)
        for b in blocks:
            txt = latex2plain((b.get("text") or "").strip())
            box = b.get("box_2d")
            if not txt or not box or len(box) != 4:
                continue
            ymin, xmin, ymax, xmax = box
            align = b.get("align", "left")
            # 표 숫자 셀은 OCR align이 center/right로 혼재 → 우측정렬로 통일(열 가지런히)
            if _is_numeric_cell(txt):
                align = "right"
            # 박스는 글자에 타이트하게 유지(폭 여유 X). 그래야 박스≈글자라 정렬 기준점이
            # 보존되고 OCR align 오차의 영향이 최소화된다.
            x = int(xmin/1000 * EMU_W); y = int(ymin/1000 * EMU_H)
            w_emu = max((xmax-xmin)/1000, 0.04) * EMU_W
            h_emu = max((ymax-ymin)/1000, 0.02) * EMU_H
            # 폰트 크기: OCR 추정이 아니라 박스 기하로 산정(모순값 자동 보정)
            fp = _fit_font(txt, w_emu/12700.0, h_emu/12700.0) * _FONT_SAFETY
            tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(int(w_emu)), Emu(int(h_emu)))
            tf = tb.text_frame
            # 긴 문단만 컬럼 폭에 맞춰 줄바꿈. 짧은 셀(숫자/라벨)은 줄바꿈 끔(약간 넘쳐도 깨지지 않게).
            tf.word_wrap = _wrap_intent(txt)
            tf.auto_size = None
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.TOP   # 위→아래로 흘러 이웃 블록 침범 방지
            col = (b.get("color") or "FFFFFF").replace("#", "")
            for j, line in enumerate(txt.split("\n")):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.alignment = ALIGN.get(align, PP_ALIGN.LEFT)
                run = p.add_run(); run.text = line
                run.font.size = Pt(round(fp, 1)); run.font.name = font
                if b.get("bold"):
                    run.font.bold = True
                try:
                    run.font.color.rgb = RGBColor.from_string(col if len(col) == 6 else "FFFFFF")
                except Exception:
                    run.font.color.rgb = RGBColor.from_string("FFFFFF")
    prs.save(out_path)
    return out_path


# ──────────────────────────────────────────────────────────────────────────
# 오케스트레이션 (전체 파이프라인)
# ──────────────────────────────────────────────────────────────────────────
def convert(input_path: str, output_path: str, workdir: str | None = None,
            erase: bool = True, font: str | None = None,
            api_key: str | None = None, max_workers: int = 6,
            tables: bool = False) -> str:
    """전체 변환. erase=False면 원본 배경 위에 텍스트만 올림(빠름, 글자 잔존).

    슬라이드별 글자제거/OCR은 max_workers 개 스레드로 병렬 실행(I/O 바운드라 효과 큼).
    max_workers=1 이면 순차. workdir 지정 시 convert.log 에 전체 로그 기록.

    tables=True 면 표를 네이티브 PowerPoint 표로 분리(objsep). 그림/차트는 인식하지 않고
    배경에 그대로 둔다. 표 검출에 슬라이드당 비전 1회 추가 호출. 기본은 글자만.
    """
    import tempfile
    t_start = time.time()
    workdir = workdir or tempfile.mkdtemp(prefix="nlm2pptx_")
    slides_dir = os.path.join(workdir, "slides")
    clean_dir = os.path.join(workdir, "clean")
    os.makedirs(clean_dir, exist_ok=True)
    fh = _add_file_log(workdir)
    log.info(f"=== convert 시작: {input_path} → {output_path} "
             f"(erase={erase}, workers={max_workers}, workdir={workdir}) ===")
    try:
        # 1) 추출
        t0 = time.time()
        slide_pngs = extract_slides(input_path, slides_dir)
        n = len(slide_pngs)
        log.info(f"[1/4] 슬라이드 추출: {n}장 ({time.time()-t0:.1f}s)")

        # 2+3) 슬라이드별 (글자제거 + OCR) 병렬
        bg_paths = [None] * n
        ocr_results = [None] * n
        errors = []

        def process(idx_png):
            i, png = idx_png
            tag = f"slide{i+1}"
            try:
                if erase:
                    clean = os.path.join(clean_dir, f"clean{i+1}.png")
                    erase_text(png, clean, api_key=api_key)
                    bg_paths[i] = clean
                else:
                    bg_paths[i] = png
                ocr_results[i] = ocr_slide(png, api_key=api_key)
                log.info(f"{tag}: 완료 (bg={'clean' if erase else 'orig'}, "
                         f"{len(ocr_results[i])}블록)")
            except Exception as e:
                log.error(f"{tag}: 처리 실패 — {type(e).__name__}: {e}")
                errors.append((i+1, str(e)))
                # 실패해도 진행: 배경은 원본, 텍스트는 빈 것으로 폴백
                bg_paths[i] = bg_paths[i] or png
                ocr_results[i] = ocr_results[i] or []

        stage = "글자제거+OCR" if erase else "OCR"
        log.info(f"[2/4][3/4] {stage} 병렬 시작 ({n}장, workers={max_workers})")
        t0 = time.time()
        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            list(ex.map(process, list(enumerate(slide_pngs))))
        log.info(f"[2/4][3/4] {stage} 완료 ({time.time()-t0:.1f}s, 실패 {len(errors)}장)")

        # 4) 조립
        t0 = time.time()
        if tables:
            import objsep
            log.info(f"[4/4] 글씨+표 조립(표 분리) → {output_path}")
            objsep.build_objsep(slide_pngs, bg_paths, ocr_results, output_path,
                                font=font, det_cache_dir=os.path.join(workdir, "objsep_det"))
        else:
            build_pptx(bg_paths, ocr_results, output_path, font=font)
        log.info(f"[4/4] PPTX 조립 → {output_path} ({time.time()-t0:.1f}s)")

        if errors:
            log.warning(f"일부 슬라이드 실패(빈 텍스트로 대체): {[e[0] for e in errors]}")
        log.info(f"=== 완료: {output_path} (총 {time.time()-t_start:.1f}s) ===")
        return output_path
    finally:
        log.removeHandler(fh)
        fh.close()


def _cli():
    ap = argparse.ArgumentParser(description="NotebookLM 이미지 PPTX/PDF → 편집 가능 PPTX")
    ap.add_argument("input", help="입력 .pptx 또는 .pdf")
    ap.add_argument("output", help="출력 .pptx")
    ap.add_argument("--workdir", default=None, help="중간 파일 디렉토리(기본: 임시)")
    ap.add_argument("--no-erase", action="store_true", help="글자 제거 생략(원본 배경 유지, 빠름)")
    ap.add_argument("--font", default=None, help=f"한글 폰트명(기본: {DEFAULT_FONT})")
    ap.add_argument("--workers", type=int, default=6, help="슬라이드 병렬 처리 수(기본 6, 1=순차)")
    ap.add_argument("--tables", action="store_true",
                    help="표를 네이티브 PowerPoint 표로 분리(글씨+표). 그림/차트는 배경에 유지. "
                         "표 검출에 슬라이드당 비전 호출 1회 추가")
    args = ap.parse_args()
    convert(args.input, args.output, workdir=args.workdir,
            erase=not args.no_erase, font=args.font, max_workers=args.workers,
            tables=args.tables)


if __name__ == "__main__":
    _cli()
