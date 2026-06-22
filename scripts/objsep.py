"""객체 분리(Object separation): 통이미지 슬라이드를 '편집 가능한 객체'로 분리해 PPTX 조립.

레이어 구성:
  (1) 배경 = 글자+객체 제거된 깨끗한 이미지(따로 생성, 원본 비율)
  (2) 표 = 네이티브 PPTX 표(셀 텍스트=OCR 매핑, 열너비·행높이·셀색을 원본에서 추정)
  (3) 그림/아이콘/차트 = 글자 지운 배경에서 크롭한 개별 이미지
  (4) 그 외 글자 = 편집 가능한 텍스트박스

검출기:
  - C: gpt-5.5 비전이 객체+표구조 JSON 반환 (설치 불필요, API 비용)
  - A: img2table(표 격자)+OpenCV(그림). 추가 API 없음 (선택 의존: img2table, opencv-python)
  - hybrid(기본): 표=A 네이티브 표, 그림=C. A가 못 잡고 C가 구조를 준 표는 C 구조 사용.

build_pptx(텍스트만)와 달리 nlm2pptx.convert(..., objects=True)에서 호출된다.
"""
import base64, json, os, sys, time, urllib.request
import numpy as np
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import nlm2pptx
from nlm2pptx import _fit_font, _wrap_intent, _is_numeric_cell, latex2plain, _em_width
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image

try:
    import cv2
except Exception:                      # opencv 없으면 배경 인페인트/그림CV 비활성(평면 채움 폴백)
    cv2 = None

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_CROP_SEQ = 0

# ───────────────────────── 검출기 C (gpt-5.5 비전) ─────────────────────────
OBJ_PROMPT = """You are a slide LAYOUT analyzer. The image is one slide. Return ONLY JSON: {"objects":[...]}.
Detect ONLY NON-TEXT visual objects, each as one block:
- type: "image" (photo/illustration/figure), "icon" (logo/badge/pictogram), "chart" (bar/line/pie graph), or "table" (data grid).
- box_2d: [ymin,xmin,ymax,xmax] in 0-1000 coords, TIGHT around the object.
For "table" also include: n_rows, n_cols (integers), and cells: [{"r":int,"c":int,"text":str,"box_2d":[...]}] (r,c 0-based, row 0 = top/header).
Do NOT output plain text, titles, bullets, or captions as objects (those are handled separately). Only real visual objects/tables.
If there are none, return {"objects":[]}. Be precise with boxes."""

def detect_C(image_path, ocr_blocks=None, model=None, retries=3, api_key=None):
    key = nlm2pptx._api_key(api_key)
    model = model or nlm2pptx.OCR_MODEL
    img_b64 = base64.b64encode(open(image_path, "rb").read()).decode()
    body = {"model": model, "messages": [{"role": "user", "content": [
        {"type": "text", "text": OBJ_PROMPT},
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}}]}],
        "response_format": {"type": "json_object"}}
    last = None
    for _ in range(retries):
        try:
            req = urllib.request.Request(f"{nlm2pptx.OPENAI_BASE}/chat/completions",
                data=json.dumps(body).encode(),
                headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"})
            d = json.loads(urllib.request.urlopen(req, timeout=200).read())
            return json.loads(d["choices"][0]["message"]["content"]).get("objects", [])
        except Exception as e:
            last = f"{type(e).__name__}: {e}"; time.sleep(4)
    raise RuntimeError(f"detect_C 실패 {retries}회: {last}")

# ───────────────────────── 검출기 A (img2table + OpenCV) ─────────────────────────
def _iou(a, b):
    ay1, ax1, ay2, ax2 = a; by1, bx1, by2, bx2 = b
    ix1, iy1, ix2, iy2 = max(ax1, bx1), max(ay1, by1), min(ax2, bx2), min(ay2, by2)
    inter = max(0, ix2-ix1) * max(0, iy2-iy1)
    if inter == 0:
        return 0.0
    ua = (ax2-ax1)*(ay2-ay1) + (bx2-bx1)*(by2-by1) - inter
    return inter/ua if ua > 0 else 0.0

def _detect_figures_cv(image_path, exclude, W, H):
    if cv2 is None:
        return []
    img = cv2.imread(image_path)
    if img is None:
        return []
    edges = cv2.Canny(cv2.cvtColor(img, cv2.COLOR_BGR2GRAY), 40, 120)
    edges = cv2.dilate(edges, np.ones((25, 25), np.uint8), iterations=2)
    cnts, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    out = []
    for cnt in cnts:
        x, y, w, h = cv2.boundingRect(cnt)
        if (w*h)/(W*H) < 0.05 or (w*h)/(W*H) > 0.85 or w < 60 or h < 60:
            continue
        box = [y/H*1000, x/W*1000, (y+h)/H*1000, (x+w)/W*1000]
        if not any(_iou(box, ex) > 0.25 for ex in exclude):
            out.append({"type": "image", "box_2d": box})
    return out

def detect_A(image_path, ocr_blocks=None):
    """img2table 표 격자 + OpenCV 그림. img2table 미설치 시 그림(CV)만."""
    W, H = Image.open(image_path).size
    objs = []
    try:
        from img2table.document import Image as I2T
        doc = I2T(image_path, detect_rotation=False)
        tables = doc.extract_tables(borderless_tables=True, implicit_rows=True,
                                    implicit_columns=True, min_confidence=30)
    except Exception:
        tables = []
    for t in tables:
        bb = t.bbox
        rows = list(t.content.values())
        nr = len(rows); nc = max((len(r) for r in rows), default=0)
        if nr < 2 or nc < 2:
            continue
        cells = [{"r": r, "c": c, "box_2d": [cell.bbox.y1/H*1000, cell.bbox.x1/W*1000,
                                             cell.bbox.y2/H*1000, cell.bbox.x2/W*1000]}
                 for r, row in enumerate(rows) for c, cell in enumerate(row)]
        objs.append({"type": "table",
                     "box_2d": [bb.y1/H*1000, bb.x1/W*1000, bb.y2/H*1000, bb.x2/W*1000],
                     "table": {"n_rows": nr, "n_cols": nc, "cells": cells}})
    objs += _detect_figures_cv(image_path, [o["box_2d"] for o in objs], W, H)
    return objs

# ───────────────────────── 하이브리드 ─────────────────────────
def _table_filled_ratio(obj, ocr_blocks):
    cells = (obj.get("table") or {}).get("cells") or []
    if not cells:
        return 0.0
    return sum(1 for c in cells if c.get("box_2d") and _ocr_text_in(c["box_2d"], ocr_blocks).strip()) / len(cells)

def _has_structure(o):
    t = o.get("table") or {}
    try:
        return int(t.get("n_rows") or 0) >= 2 and int(t.get("n_cols") or 0) >= 2 and bool(t.get("cells"))
    except Exception:
        return False

def combine_hybrid(a_objs, c_objs, ocr_blocks, min_fill=0.3):
    tables = [o for o in a_objs
              if o.get("type") == "table" and _table_filled_ratio(o, ocr_blocks) >= min_fill]
    figs = []
    for o in c_objs or []:
        box = o.get("box_2d")
        if not box or len(box) != 4:
            continue
        over = any(_iou(box, tb["box_2d"]) > 0.3 or _center_in(box, tb["box_2d"]) for tb in tables)
        if o.get("type") == "table":
            if _has_structure(o) and not over:
                tables.append(o); continue
            o = dict(o); o["type"] = "image"
        if over:
            continue
        figs.append(o)
    return tables + figs

def detect_hybrid(image_path, ocr_blocks=None):
    return combine_hybrid(detect_A(image_path, ocr_blocks), detect_C(image_path, ocr_blocks), ocr_blocks)

_DETECTORS = {"hybrid": detect_hybrid, "C": detect_C, "A": detect_A}

# ───────────────────────── 공통 Assembler ─────────────────────────
def _center_in(box, region):
    cy = (box[0] + box[2]) / 2; cx = (box[1] + box[3]) / 2
    return region[0] <= cy <= region[2] and region[1] <= cx <= region[3]

def _ocr_text_in(region, ocr_blocks):
    hits = [(b["box_2d"][0], b["box_2d"][1], (b.get("text") or "").strip())
            for b in (ocr_blocks or [])
            if b.get("box_2d") and len(b["box_2d"]) == 4 and _center_in(b["box_2d"], region)]
    hits.sort()
    return " ".join(t for _, _, t in hits if t)

def _ocr_cell(region, ocr_blocks):
    hits = [b for b in (ocr_blocks or [])
            if b.get("box_2d") and len(b["box_2d"]) == 4 and _center_in(b["box_2d"], region)]
    hits.sort(key=lambda b: (b["box_2d"][0], b["box_2d"][1]))
    text = " ".join((b.get("text") or "").strip() for b in hits if (b.get("text") or "").strip())
    color = next((b.get("color") for b in hits if b.get("color")), None)
    align = hits[0].get("align") if hits else None
    bold = any(b.get("bold") for b in hits)
    return text, color, align, bold

def _cell_font(text, w_pt, h_pt):
    lines = (text or " ").split("\n")
    nl = max(1, len(lines))
    longest = max((_em_width(l) for l in lines), default=1.0) or 1.0
    total = sum(_em_width(l) for l in lines) or 1.0
    f_h = h_pt / nl
    f = f_h if longest * f_h <= w_pt else (w_pt * h_pt / (1.4 * total)) ** 0.5
    return max(5.0, min(f, 13.0))

def _table_valid(obj):
    t = obj.get("table") or {}
    try:
        return int(t.get("n_rows") or 0) >= 1 and int(t.get("n_cols") or 0) >= 1
    except Exception:
        return False

def _median_color(orig_np, box, OW, OH):
    if orig_np is None or not box:
        return None
    ymin, xmin, ymax, xmax = box
    x1, y1 = max(0, int(xmin/1000*OW)), max(0, int(ymin/1000*OH))
    x2, y2 = min(OW, int(xmax/1000*OW)), min(OH, int(ymax/1000*OH))
    if x2 <= x1 or y2 <= y1:
        return None
    med = np.median(orig_np[y1:y2, x1:x2].reshape(-1, 3), axis=0)
    return tuple(int(v) for v in med)

def _set_no_grid_style(tbl):
    from pptx.oxml.ns import qn
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is None:
        return
    for el in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(el)
    sid = tblPr.makeelement(qn('a:tableStyleId'), {})
    sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'   # No Style, No Grid
    tblPr.append(sid)

def _add_table(slide, obj, font, ocr_blocks=None, orig_np=None, OW=0, OH=0):
    box = obj["box_2d"]; t = obj.get("table") or {}
    nr, nc = int(t.get("n_rows") or 0), int(t.get("n_cols") or 0)
    cells = t.get("cells") or []
    if nr < 1 or nc < 1:
        return False
    ymin, xmin, ymax, xmax = box
    x = int(xmin/1000*EMU_W); y = int(ymin/1000*EMU_H)
    w = int(max((xmax-xmin)/1000, 0.05)*EMU_W); h = int(max((ymax-ymin)/1000, 0.05)*EMU_H)
    tbl = slide.shapes.add_table(nr, nc, Emu(x), Emu(y), Emu(w), Emu(h)).table
    _set_no_grid_style(tbl)
    tbl.first_row = tbl.last_row = tbl.first_col = tbl.last_col = False
    tbl.horz_banding = tbl.vert_banding = False
    # 열너비·행높이 = 셀 박스의 열별/행별 중앙값(병합 헤더 소수는 median이 무시 → 데이터 격자 폭)
    colmin, colmax, rowmin, rowmax = {}, {}, {}, {}
    for cell in cells:
        cb = cell.get("box_2d")
        if not cb or len(cb) != 4:
            continue
        cc, rr = int(cell.get("c", -1)), int(cell.get("r", -1))
        if 0 <= cc < nc:
            colmin.setdefault(cc, []).append(cb[1]); colmax.setdefault(cc, []).append(cb[3])
        if 0 <= rr < nr:
            rowmin.setdefault(rr, []).append(cb[0]); rowmax.setdefault(rr, []).append(cb[2])
    for cc in range(nc):
        if cc in colmin:
            wc = (float(np.median(colmax[cc])) - float(np.median(colmin[cc]))) / 1000 * EMU_W
            if wc > 0:
                tbl.columns[cc].width = Emu(int(wc))
    for rr in range(nr):
        if rr in rowmin:
            hr = (float(np.median(rowmax[rr])) - float(np.median(rowmin[rr]))) / 1000 * EMU_H
            if hr > 0:
                tbl.rows[rr].height = Emu(int(hr))
    cell_w0 = (w/12700.0)/max(nc, 1); cell_h0 = (h/12700.0)/max(nr, 1)
    for cell in cells:
        r, c = int(cell.get("r", -1)), int(cell.get("c", -1))
        if not (0 <= r < nr and 0 <= c < nc):
            continue
        cbox = cell.get("box_2d")
        text = (cell.get("text") or "").strip()
        color = align = None; bold = False
        if cbox:
            ot, color, align, bold = _ocr_cell(cbox, ocr_blocks)
            if not text:
                text = ot
        tc = tbl.cell(r, c)
        tc.margin_left = tc.margin_right = Emu(27432)
        tc.margin_top = tc.margin_bottom = Emu(9144)
        tc.vertical_anchor = MSO_ANCHOR.MIDDLE
        fillc = _median_color(orig_np, cbox, OW, OH)
        if fillc:
            tc.fill.solid(); tc.fill.fore_color.rgb = RGBColor(*fillc)
        else:
            tc.fill.background()
        tc.text = latex2plain(text)
        if color and len(color.replace('#', '')) == 6:
            txt_rgb = RGBColor.from_string(color.replace('#', ''))
        elif fillc:
            lum = 0.299*fillc[0] + 0.587*fillc[1] + 0.114*fillc[2]
            txt_rgb = RGBColor(0x22, 0x22, 0x22) if lum > 140 else RGBColor(0xFF, 0xFF, 0xFF)
        else:
            txt_rgb = RGBColor(0x22, 0x22, 0x22)
        al = PP_ALIGN.RIGHT if _is_numeric_cell(text) else \
            {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        if cbox:
            cw = max((cbox[3]-cbox[1])/1000, 0.02) * EMU_W / 12700.0
            ch = max((cbox[2]-cbox[0])/1000, 0.02) * EMU_H / 12700.0
        else:
            cw, ch = cell_w0, cell_h0
        fp = _cell_font(latex2plain(text), max(cw - 4.0, 6.0), ch)
        for para in tc.text_frame.paragraphs:
            para.alignment = al
            for run in para.runs:
                run.font.size = Pt(round(fp, 1)); run.font.name = font
                run.font.color.rgb = txt_rgb
                if bold:
                    run.font.bold = True
    return True

def _place_text(slide, txt, box, align, bold, color, font):
    txt = latex2plain((txt or "").strip())
    if not txt:
        return
    if _is_numeric_cell(txt):
        align = "right"
    ymin, xmin, ymax, xmax = box
    x = int(xmin/1000*EMU_W); y = int(ymin/1000*EMU_H)
    w = max((xmax-xmin)/1000, 0.04) * EMU_W
    h = max((ymax-ymin)/1000, 0.02) * EMU_H
    fp = _fit_font(txt, w/12700.0, h/12700.0) * 0.98
    tf = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(int(w)), Emu(int(h))).text_frame
    tf.word_wrap = _wrap_intent(txt); tf.auto_size = None
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    col = (color or "FFFFFF").replace("#", "")
    for j, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = ALIGN.get(align, PP_ALIGN.LEFT)
        r = p.add_run(); r.text = line
        r.font.size = Pt(round(fp, 1)); r.font.name = font
        if bold:
            r.font.bold = True
        try:
            r.font.color.rgb = RGBColor.from_string(col if len(col) == 6 else "FFFFFF")
        except Exception:
            r.font.color.rgb = RGBColor.from_string("FFFFFF")

def _build_background(bg_path, regions, tmpdir):
    """배경에서 객체 영역 제거(작은 영역 인페인트, 큰 영역 배경색 채움). opencv 없으면 원본 배경 사용."""
    global _CROP_SEQ
    if cv2 is None or not regions:
        return bg_path
    img = cv2.imread(bg_path)
    if img is None:
        return bg_path
    H, W = img.shape[:2]; b = 8
    border = np.concatenate([img[:b].reshape(-1, 3), img[-b:].reshape(-1, 3),
                             img[:, :b].reshape(-1, 3), img[:, -b:].reshape(-1, 3)])
    col = np.median(border, axis=0)
    mask = np.zeros((H, W), np.uint8)
    for (ymin, xmin, ymax, xmax) in regions:
        x1, y1 = max(0, int(xmin/1000*W)), max(0, int(ymin/1000*H))
        x2, y2 = min(W, int(xmax/1000*W)), min(H, int(ymax/1000*H))
        if x2 <= x1 or y2 <= y1:
            continue
        if (x2-x1)*(y2-y1) > 0.18*W*H:
            img[y1:y2, x1:x2] = col
        else:
            mask[y1:y2, x1:x2] = 255
    if mask.any():
        img = cv2.inpaint(img, mask, 6, cv2.INPAINT_TELEA)
    _CROP_SEQ += 1
    out = os.path.join(tmpdir, f"background_{_CROP_SEQ}.png")
    cv2.imwrite(out, img)
    return out

def _assemble_into(slide, bg_path, original_path, ocr_blocks, objects, font, tmpdir):
    global _CROP_SEQ
    orig = Image.open(original_path).convert("RGB")
    OW, OH = orig.size
    orig_np = np.array(orig)
    clean_img = Image.open(bg_path).convert("RGB")
    if clean_img.size != (OW, OH):       # erase 결과 비율이 다르면 원본 비율로 복원
        clean_img = clean_img.resize((OW, OH))
        _CROP_SEQ += 1
        bg_path = os.path.join(tmpdir, f"clean_fixed_{_CROP_SEQ}.png")
        clean_img.save(bg_path)
    cW, cH = clean_img.size
    valid = [o for o in objects if o.get("box_2d") and len(o["box_2d"]) == 4]
    fill_regions = [o["box_2d"] for o in valid]
    table_regions = [o["box_2d"] for o in valid if o.get("type") == "table" and _table_valid(o)]

    slide.shapes.add_picture(_build_background(bg_path, fill_regions, tmpdir), 0, 0, width=EMU_W, height=EMU_H)

    for obj in valid:
        box = obj["box_2d"]
        if obj.get("type") == "table" and _table_valid(obj):
            _add_table(slide, obj, font, ocr_blocks, orig_np, OW, OH)
            continue
        ymin, xmin, ymax, xmax = box
        L, T, R, B = int(xmin/1000*cW), int(ymin/1000*cH), int(xmax/1000*cW), int(ymax/1000*cH)
        if R-L < 3 or B-T < 3:
            continue
        crop = clean_img.crop((max(0, L), max(0, T), min(cW, R), min(cH, B)))
        _CROP_SEQ += 1
        cp = os.path.join(tmpdir, f"obj_{_CROP_SEQ}.png"); crop.save(cp)
        slide.shapes.add_picture(cp, Emu(int(xmin/1000*EMU_W)), Emu(int(ymin/1000*EMU_H)),
                                 width=Emu(int((xmax-xmin)/1000*EMU_W)), height=Emu(int((ymax-ymin)/1000*EMU_H)))

    for b in ocr_blocks:
        box = b.get("box_2d")
        if not box or len(box) != 4 or any(_center_in(box, tr) for tr in table_regions):
            continue
        _place_text(slide, b.get("text"), box, b.get("align", "left"), b.get("bold"), b.get("color"), font)

# ───────────────────────── 멀티슬라이드 빌드 ─────────────────────────
def build_objsep(slide_pngs, bg_paths, ocr_per_slide, out_path, font=None,
                 detector="hybrid", det_cache_dir=None):
    """원본 PNG + 글자지운 배경 + OCR블록 → 객체분리 PPTX. 슬라이드 크기는 원본 이미지 비율."""
    global EMU_W, EMU_H
    font = font or nlm2pptx.DEFAULT_FONT
    det = _DETECTORS.get(detector, detect_hybrid)
    ow, oh = Image.open(slide_pngs[0]).size
    EMU_W = Inches(13.333); EMU_H = Inches(round(13.333 * oh / ow, 4))
    if det_cache_dir:
        os.makedirs(det_cache_dir, exist_ok=True)
    tmpdir = os.path.join(os.path.dirname(os.path.abspath(out_path)) or ".", "_objsep_tmp")
    os.makedirs(tmpdir, exist_ok=True)
    prs = Presentation(); prs.slide_width, prs.slide_height = EMU_W, EMU_H
    for i, (orig, bg, blocks) in enumerate(zip(slide_pngs, bg_paths, ocr_per_slide)):
        dc = os.path.join(det_cache_dir, f"slide{i+1}.json") if det_cache_dir else None
        if dc and os.path.exists(dc):
            objs = json.load(open(dc, encoding="utf-8"))
        else:
            objs = det(orig, blocks)
            if dc:
                json.dump(objs, open(dc, "w", encoding="utf-8"), ensure_ascii=False)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _assemble_into(slide, bg, orig, blocks, objs, font, tmpdir)
    prs.save(out_path)
    return out_path
